from pptx import Presentation

# Load the presentation
prs = Presentation(r"C:\Users\ADMIN\Downloads\FYP_Presentation_Updated.pptx")

# Slide 3: Background & Problem Statement
# We want to find the shape with the text "OUR SOLUTION" and append our points, 
# or just rewrite the whole text frame if it contains it.
for shape in prs.slides[2].shapes:
    if hasattr(shape, "text") and "OUR SOLUTION" in shape.text:
        shape.text = (
            "OUR SOLUTION\n"
            " AI-Powered Career Advisory Agent specifically designed for tech professionals and students.\n"
            " Instant Resume Parsing: Eliminates manual data entry by extracting core skills directly from PDFs.\n"
            " Live Market Gap Analysis: Integrates real-time job market data (Tavily Search API) to find exact missing skills.\n"
            " Multi-Agent Architecture (LangGraph): Strategist Agent drafts roadmaps, Critic Agent iteratively refines them."
        )

# Slide 18: Integration Details
for shape in prs.slides[17].shapes:
    if hasattr(shape, "text") and "Integration Details" in shape.text:
        shape.text = (
            "7.6 Integration Details\n"
            " Google Gemini 2.5 Flash: Core LLM used for resume parsing and gap analysis.\n"
            " LangGraph: Orchestrates the Multi-Agent workflow (Strategist & Critic nodes).\n"
            " FastAPI & Uvicorn: Backend REST API framework.\n"
            " Clerk Auth: Secure, production-ready user authentication.\n"
            " React 19 + Tailwind CSS: Frontend with react-markdown for structured roadmaps.\n"
            " Tavily Search API: Retrieves live job market demands."
        )

# Slide 20: Phase 2
for shape in prs.slides[19].shapes:
    if hasattr(shape, "text") and "FR4" in shape.text:
        shape.text = (
            "PHASE 2: Automated Parsing & Multi-Agent Implementation\n"
            " Automated Resume Parsing: Replaced manual surveys with seamless PDF uploads. Backend extracts profile via Gemini.\n"
            " Live Job Market Integration: Tavily Search dynamically finds 'missing skills' based on the user's extracted profile.\n"
            " Agentic Roadmap Generation: Built a customized 3-month action plan generator using LangGraph.\n"
            " Multi-Agent Feedback Loop: Strategist Agent builds the plan; Critic Agent reviews and requests revisions."
        )

# Slide 24: Conclusion
for shape in prs.slides[23].shapes:
    if hasattr(shape, "text") and "What Was Achieved" in shape.text:
        shape.text = (
            "What Was Achieved\n"
            " Transitioned from manual forms to an Automated Agentic Workflow.\n"
            " Validated that LangGraph can orchestrate multiple AI agents (Strategist + Critic) to iteratively improve roadmaps.\n"
            " Proved that grounding LLMs in Live Web Search (Tavily) prevents hallucinations and provides accurate skill gaps.\n"
            " Deployed a secure, modern frontend with Clerk Authentication and robust PDF handling."
        )

# Slide 25: Future Scope
for shape in prs.slides[24].shapes:
    if hasattr(shape, "text") and "Phase 2 Enhancements" in shape.text:
        shape.text = (
            "Future Scope & Next Steps\n"
            " Long-Term Episodic Memory: Allow the AI to remember user progress across sessions and track roadmap completion.\n"
            " Interactive Interviewer Agent: A sub-agent that conducts mock interviews based on skills identified in the gap analysis.\n"
            " Algorithmic Mentor Matching: Connect users with senior IT professionals using vector embeddings (ChromaDB).\n"
            " Automated Job Auto-Apply Agent: An autonomous browser agent that drafts cover letters and applies to matching jobs."
        )

# Save the updated presentation
output_path = r"C:\Users\ADMIN\Downloads\FYP_Presentation_Agentic_Updated.pptx"
prs.save(output_path)
print(f"Saved successfully to {output_path}")
